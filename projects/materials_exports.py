"""Phase 2 — Builders for AI-generated teaching materials.

Three builders:
  • build_challenge_card_pptx(project, week, content)  — template-based 4-slide PPT
  • build_lesson_plan_docx(project, week, content)     — from-scratch DOCX (~16 pages)
  • build_session_pptx(project, session, content)      — from-scratch PPT (~17-21 slides)

Each takes structured JSON `content` (produced by AI) so the file can be
rebuilt without re-calling the AI.
"""
import copy
import io
import re
from pathlib import Path

from django.conf import settings

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Pt, Cm, RGBColor
from docx.oxml import OxmlElement
from docx.oxml.ns import qn

from pptx import Presentation
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from pptx.util import Inches, Pt as PPt, Emu

try:
    from PIL import Image as _PILImage
    _PIL_OK = True
except Exception:
    _PIL_OK = False

from .stock_images import fetch_stock_image_io, prefetch_images

from framework.models import Week, Session
from .models import Project
from .docx_utils import (
    PRIMARY, INK, MUTED, ACCENT_TEAL, ACCENT_AMBER, ACCENT_PINK, ACCENT_BLUE,
    set_cell_shading, add_styled_run, render_inline, render_markdown,
)


CHALLENGE_TEMPLATE = settings.BASE_DIR / 'Materials' / 'Challeneg Cards' / 'ChallengeCard_Template 1.pptx'


# ╔══════════════════════════════════════════════════════════════════════╗
# ║  CHALLENGE CARD — template-based                                     ║
# ╚══════════════════════════════════════════════════════════════════════╝

# Defensive word caps — applied even if AI exceeds the prompt's word limits.
CC_LIMITS = {
    'big_question':   16,
    'job_interest':   13,
    'scenario':       55,
    'task':           28,
    'level_up':       32,
    'guideline':      15,
    'word':           1,    # single word
    'think_about':    22,
    'i_wonder':       28,
    'portfolio':      40,
}


def _cap_words(text: str, max_words: int) -> str:
    """Hard-truncate to N words; preserves trailing punctuation if room allows."""
    if not text:
        return ''
    s = str(text).strip()
    parts = s.split()
    if len(parts) <= max_words:
        return s
    truncated = ' '.join(parts[:max_words]).rstrip(' ,;:-—')
    # Re-attach a sensible terminator if we cut off mid-sentence
    if not truncated.endswith(('.', '?', '!', '…')):
        truncated += '…'
    return truncated


def _word_capped_lines(items, max_words: int) -> str:
    return '\n'.join(_cap_words(str(x), max_words) for x in items if str(x).strip())


def _replace_run_text(text_frame, new_text: str, *, lock_size: bool = True):
    """Replace the entire text of a TextFrame, keeping the formatting of the
    first run. Newline-separated lines become separate paragraphs."""
    if not new_text:
        new_text = ''

    # Capture font from first run as the template
    first_para = text_frame.paragraphs[0] if text_frame.paragraphs else None
    template_run = None
    if first_para and first_para.runs:
        template_run = first_para.runs[0]

    # Clear all paragraphs
    p_elem = text_frame._txBody
    paragraphs = list(p_elem.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}p'))
    for p in paragraphs[1:]:
        p_elem.remove(p)

    first_p = text_frame.paragraphs[0]
    for run in list(first_p.runs):
        run._r.getparent().remove(run._r)

    lines = str(new_text).split('\n')
    for idx, line in enumerate(lines):
        if idx == 0:
            para = first_p
        else:
            para = text_frame.add_paragraph()
        run = para.add_run()
        run.text = line
        if template_run is not None:
            tr_font = template_run.font
            try:
                if tr_font.size:
                    run.font.size = tr_font.size
            except Exception:
                pass
            try:
                if tr_font.name:
                    run.font.name = tr_font.name
            except Exception:
                pass
            try:
                if tr_font.bold is not None:
                    run.font.bold = tr_font.bold
            except Exception:
                pass
            try:
                if tr_font.color and tr_font.color.type is not None:
                    run.font.color.rgb = tr_font.color.rgb
            except Exception:
                pass

    # Lock auto-size so neither shape grows nor text shrinks — protects layout
    if lock_size:
        try:
            text_frame.auto_size = None
            text_frame.word_wrap = True
        except Exception:
            pass


def _bullet_lines(items):
    return '\n'.join(f'• {x}' for x in items if str(x).strip())


def _replace_shape_with_image(slide, shape, query: str, *, variant: int = 0):
    """Drop a stock image at the position/size of `shape` and remove the
    placeholder text. Used for [PROJECT GRAPHIC / ICON] slots on slide 1."""
    try:
        x, y, w, h = shape.left, shape.top, shape.width, shape.height
        # Hide the placeholder text
        if shape.has_text_frame:
            for p in list(shape.text_frame._txBody.findall(
                '{http://schemas.openxmlformats.org/drawingml/2006/main}p')):
                shape.text_frame._txBody.remove(p)
            # Re-add an empty paragraph so the shape stays valid
            from pptx.oxml import parse_xml as _px
            shape.text_frame._txBody.append(_px(
                '<a:p xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
            ))
        img = fetch_stock_image_io(query, orientation='squarish',
                                   variant=variant)
        if img:
            slide.shapes.add_picture(img, x, y, width=w, height=h)
            return True
    except Exception:
        pass
    return False


def build_challenge_card_pptx(project: Project, week: Week, content: dict) -> io.BytesIO:
    """Build a MODERN 4-slide Challenge Card from scratch (no template).

    Design language matches the Session PPT: clean typography, generous
    whitespace, embedded stock images with overlays, brand palette.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    sw, sh = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    raw_grade   = str(project.grade or '').strip()
    grade_clean = re.sub(r'^grade\s*', '', raw_grade, flags=re.IGNORECASE).strip()
    grade       = f'Grade {grade_clean}' if grade_clean else raw_grade
    work_form   = project.get_subject_track_display()
    challenge_n = week.number

    # Apply hard word caps so even verbose AI output stays clean.
    big_qs       = [_cap_words(x, CC_LIMITS['big_question']) for x in (content.get('big_questions') or [])][:3]
    scenario     = _cap_words(content.get('scenario', ''), CC_LIMITS['scenario'])
    raw_tasks    = content.get('tasks') or []
    tasks        = [_cap_words(t, CC_LIMITS['task']) for t in raw_tasks[:3]]
    if len(raw_tasks) > 3:
        tasks.append(_cap_words(raw_tasks[3], CC_LIMITS['level_up']))
    guidelines   = [_cap_words(g, CC_LIMITS['guideline']) for g in (content.get('guidelines') or [])][:5]
    words_raw    = content.get('words_of_week') or []
    words        = []
    for w in words_raw[:3]:
        first = re.sub(r'[^A-Za-z]', '', str(w).split()[0]) if str(w).split() else ''
        words.append(first or str(w).strip())
    think_about  = [_cap_words(t, CC_LIMITS['think_about']) for t in (content.get('think_about') or [])][:3]
    i_wonder     = _cap_words(content.get('i_wonder', ''), CC_LIMITS['i_wonder'])
    portfolio    = _cap_words(content.get('portfolio_task', ''), CC_LIMITS['portfolio'])

    # ─────────────────────────────────────────────────────────────────────
    # SLIDE 1 — Cover with Big Questions
    # ─────────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank)

    hero_query = f'{project.topic} indian school students'
    has_img = _add_image_full_bleed(slide, sw, sh, hero_query,
                                    fallback_color=SLIDE_DEEP, variant=20 + challenge_n)
    if has_img:
        _add_dark_overlay(slide, Inches(0), Inches(0), sw, sh,
                          alpha=78, top_hex='1E1B4B', bottom_hex='000000')
    else:
        _add_rect(slide, Inches(0), Inches(0), sw, sh, fill=SLIDE_DEEP,
                  shape=MSO_SHAPE.RECTANGLE)

    # Top bar: project title + grade pill
    _add_textbox(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.4),
                 'NEORISE FSL  ·  CHALLENGE CARD',
                 font_size=11, bold=True, color=ACCENT_GOLD)
    _add_textbox(slide, Inches(0.8), Inches(0.95), Inches(9), Inches(0.5),
                 project.topic[:80], font_size=15, bold=True, color=WHITE)

    # Top-right pill cluster
    pill_w = Inches(2.0)
    px = sw - pill_w - Inches(0.7)
    pill_grade = _add_rect(slide, px, Inches(0.55), pill_w, Inches(0.55),
                           fill=ACCENT_GOLD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = pill_grade.text_frame
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = f'{grade}'
    r.font.size = PPt(15)
    r.font.bold = True
    r.font.color.rgb = SLIDE_DEEP

    # Challenge number — large outlined
    _add_textbox(slide, sw - Inches(2.7), Inches(1.25), Inches(2.0),
                 Inches(0.4), f'CHALLENGE  ·  #{challenge_n}',
                 font_size=12, bold=True, color=ACCENT_GOLD,
                 align=PP_ALIGN.RIGHT)

    # Eyebrow
    _add_textbox(slide, Inches(0.8), Inches(2.7), Inches(11), Inches(0.4),
                 'BIG QUESTION', font_size=14, bold=True, color=ACCENT_GOLD)

    # Big headline question (first big_q is the hero)
    hero_q = big_qs[0] if big_qs else 'What problem are we solving this week?'
    _add_textbox(slide, Inches(0.8), Inches(3.2), Inches(11.7),
                 Inches(2.0), hero_q, font_size=42, bold=True,
                 color=WHITE, parse_md=True)

    # Supporting questions row
    if len(big_qs) > 1:
        sub_y = Inches(5.6)
        gap   = Inches(0.25)
        avail = sw - Inches(1.6) - gap * (len(big_qs) - 2)
        card_w = avail / max(1, len(big_qs) - 1)
        x = Inches(0.8)
        for q in big_qs[1:3]:
            card = _add_rect(slide, x, sub_y, card_w, Inches(1.3),
                             fill=PRGBColor(0xFF, 0xFF, 0xFF),
                             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            _set_solid_fill_alpha(card, 'FFFFFF', 12)
            tf = card.text_frame
            tf.margin_left = Inches(0.3)
            tf.margin_top = Inches(0.25)
            tf.margin_right = Inches(0.3)
            tf.word_wrap = True
            _render_inline_pptx(tf.paragraphs[0], q, font_size=15,
                                color=WHITE)
            x += card_w + gap

    # Footer
    _add_textbox(slide, Inches(0.8), sh - Inches(0.55), Inches(11.7),
                 Inches(0.35), f'Week {challenge_n}  ·  {work_form}',
                 font_size=11, color=PRGBColor(0xCB, 0xD5, 0xE1))

    # ─────────────────────────────────────────────────────────────────────
    # SLIDE 2 — Scenario + Tasks
    # ─────────────────────────────────────────────────────────────────────
    slide = _new_slide(prs, blank)

    # Top accent bar
    _add_rect(slide, Inches(0), Inches(0), sw, Inches(0.18),
              fill=SLIDE_PRIMARY, shape=MSO_SHAPE.RECTANGLE)

    _add_textbox(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.35),
                 f'CHALLENGE #{challenge_n}  ·  THE SCENARIO',
                 font_size=11, bold=True, color=SLIDE_PRIMARY)
    _add_textbox(slide, Inches(0.8), Inches(0.7), Inches(12), Inches(0.6),
                 'Here\'s what\'s happening', font_size=24, bold=True,
                 color=SLIDE_INK)
    _add_rect(slide, Inches(0.8), Inches(1.35), Inches(0.6), Inches(0.06),
              fill=ACCENT_GOLD, shape=MSO_SHAPE.RECTANGLE)

    # Scenario block — image on left, scenario text on right (compressed)
    img_added = _add_image_box(slide, Inches(0.8), Inches(1.65), Inches(4.5),
                               Inches(2.1), f'{project.topic} india',
                               fallback_color=SLIDE_DEEP, variant=21 + challenge_n)
    sc_x = Inches(5.5) if img_added else Inches(0.8)
    sc_w = Inches(7.2) if img_added else Inches(11.7)

    scen_card = _add_rect(slide, sc_x, Inches(1.65), sc_w, Inches(2.1),
                          fill=PRGBColor(0xF1, 0xF5, 0xF9),
                          shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = scen_card.text_frame
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.2)
    tf.word_wrap = True
    try:
        from pptx.enum.text import MSO_AUTO_SIZE
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_SHAPE
    except Exception:
        pass
    _render_inline_pptx(tf.paragraphs[0], f'"{scenario}"',
                        font_size=13, color=SLIDE_INK, italic=True)

    # Tasks header
    _add_textbox(slide, Inches(0.8), Inches(3.95), Inches(12), Inches(0.35),
                 'YOUR TASKS', font_size=11, bold=True, color=SLIDE_PRIMARY)
    _add_textbox(slide, Inches(0.8), Inches(4.3), Inches(12), Inches(0.4),
                 'Three to deliver  ·  one bonus to push further',
                 font_size=13, color=SLIDE_MUTED)

    # 4-task grid (3 regular + 1 LEVEL UP) — taller cards
    task_y = Inches(4.85)
    task_h = Inches(2.45)
    task_w = (sw - Inches(1.6) - Inches(0.3) * 3) / 4
    task_x = Inches(0.8)
    palette = [ACTIVITY_THEMES[0]['main'], ACTIVITY_THEMES[1]['main'],
               ACTIVITY_THEMES[2]['main'], ACCENT_GOLD]
    for i in range(4):
        col = palette[i]
        is_bonus = (i == 3)
        # Card background
        card = _add_rect(slide, task_x, task_y, task_w, task_h,
                         fill=WHITE, line=col, line_width_pt=2,
                         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        # Number/badge top-left
        badge = _add_rect(slide, task_x + Inches(0.2), task_y + Inches(0.18),
                          Inches(0.5), Inches(0.5),
                          fill=col, shape=MSO_SHAPE.OVAL)
        tf = badge.text_frame
        tf.margin_top = Inches(0.05)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = '★' if is_bonus else str(i + 1)
        r.font.size = PPt(16)
        r.font.bold = True
        r.font.color.rgb = WHITE
        # Bonus label
        if is_bonus:
            _add_textbox(slide, task_x + Inches(0.85), task_y + Inches(0.22),
                         task_w - Inches(1.0), Inches(0.4),
                         'LEVEL UP', font_size=10, bold=True, color=col)
        # Task text — bigger box, smaller font, auto-shrink if overflow
        raw_task = tasks[i] if i < len(tasks) else '—'
        # Hard cap: trim absurdly long tasks at word boundary
        if len(raw_task) > 220:
            words = raw_task.split()
            trimmed, total = [], 0
            for w in words:
                if total + len(w) + 1 > 215:
                    break
                trimmed.append(w)
                total += len(w) + 1
            raw_task = ' '.join(trimmed).rstrip(',.;:') + '...'
        text_box = _add_textbox(
            slide, task_x + Inches(0.2), task_y + Inches(0.82),
            task_w - Inches(0.4), task_h - Inches(0.95),
            raw_task,
            font_size=11, color=SLIDE_INK, parse_md=True,
        )
        text_box.text_frame.word_wrap = True
        try:
            from pptx.enum.text import MSO_AUTO_SIZE
            text_box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_SHAPE
        except Exception:
            pass
        task_x += task_w + Inches(0.3)

    # ─────────────────────────────────────────────────────────────────────
    # SLIDE 3 — Guidelines + Words of the Week + Think About
    # ─────────────────────────────────────────────────────────────────────
    slide = _new_slide(prs, blank)

    _add_rect(slide, Inches(0), Inches(0), sw, Inches(0.18),
              fill=SLIDE_PRIMARY, shape=MSO_SHAPE.RECTANGLE)

    _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.4),
                 'CHALLENGE TOOLKIT', font_size=11, bold=True,
                 color=SLIDE_PRIMARY)
    _add_textbox(slide, Inches(0.8), Inches(0.85), Inches(12), Inches(0.7),
                 'Guides, words, and prompts', font_size=28, bold=True,
                 color=SLIDE_INK)
    _add_rect(slide, Inches(0.8), Inches(1.65), Inches(0.6), Inches(0.06),
              fill=ACCENT_GOLD, shape=MSO_SHAPE.RECTANGLE)

    # Two-column layout: left = Guidelines, right = (Words + Think About stacked)
    col_left_x  = Inches(0.8)
    col_left_w  = Inches(6.5)
    col_right_x = Inches(7.6)
    col_right_w = Inches(5.0)

    # ── Left column — Guidelines ──
    _add_textbox(slide, col_left_x, Inches(2.05), col_left_w, Inches(0.4),
                 '✓  GUIDELINES', font_size=12, bold=True,
                 color=ACTIVITY_THEMES[1]['main'])
    g_y = Inches(2.55)
    for i, g in enumerate(guidelines[:5]):
        # Check icon dot
        _add_decorative_dot(slide, col_left_x + Inches(0.05),
                            g_y + Inches(0.18), Inches(0.18),
                            ACTIVITY_THEMES[1]['main'])
        # Text
        gbox = _add_textbox(slide, col_left_x + Inches(0.45), g_y,
                            col_left_w - Inches(0.5), Inches(0.7), '',
                            font_size=14, color=SLIDE_INK, parse_md=True)
        _render_inline_pptx(gbox.text_frame.paragraphs[0], g, font_size=14,
                            color=SLIDE_INK)
        if i < min(len(guidelines), 5) - 1:
            _add_rect(slide, col_left_x + Inches(0.45),
                      g_y + Inches(0.7), col_left_w - Inches(0.5),
                      Inches(0.012), fill=SOFT_LINE,
                      shape=MSO_SHAPE.RECTANGLE)
        g_y += Inches(0.78)

    # ── Right column — Words of the Week ──
    _add_textbox(slide, col_right_x, Inches(2.05), col_right_w, Inches(0.4),
                 '📚  WORDS OF THE WEEK', font_size=12, bold=True,
                 color=ACTIVITY_THEMES[2]['main'])

    # Three word chips in a row
    chip_total_w = col_right_w - Inches(0.4)
    chip_w = (chip_total_w - Inches(0.3)) / 3
    cx = col_right_x
    cy = Inches(2.55)
    word_palette = [ACTIVITY_THEMES[2], ACTIVITY_THEMES[3], ACTIVITY_THEMES[4]]
    for i, w in enumerate(words[:3]):
        theme_w = word_palette[i]
        chip = _add_rect(slide, cx, cy, chip_w, Inches(0.95),
                         fill=theme_w['tint'], line=theme_w['main'],
                         line_width_pt=1.5,
                         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tf = chip.text_frame
        tf.margin_top = Inches(0.22)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = w
        r.font.size = PPt(16)
        r.font.bold = True
        r.font.color.rgb = theme_w['ink']
        cx += chip_w + Inches(0.15)

    _add_textbox(slide, col_right_x, Inches(3.65), col_right_w, Inches(0.4),
                 f'Find their meaning by end of Week {challenge_n}',
                 font_size=10, italic=True, color=SLIDE_MUTED)

    # ── Right column (continued) — Think About ──
    _add_textbox(slide, col_right_x, Inches(4.3), col_right_w, Inches(0.4),
                 '💭  THINK ABOUT', font_size=12, bold=True,
                 color=ACTIVITY_THEMES[3]['main'])
    t_y = Inches(4.8)
    for i, t in enumerate(think_about[:3]):
        _add_textbox(slide, col_right_x, t_y, Inches(0.5), Inches(0.5),
                     str(i + 1), font_size=22, bold=True,
                     color=ACTIVITY_THEMES[3]['main'])
        tbox = _add_textbox(slide, col_right_x + Inches(0.5), t_y + Inches(0.05),
                            col_right_w - Inches(0.5), Inches(0.7), '',
                            font_size=12, color=SLIDE_INK, parse_md=True)
        _render_inline_pptx(tbox.text_frame.paragraphs[0], t, font_size=12,
                            color=SLIDE_INK)
        t_y += Inches(0.72)

    # ─────────────────────────────────────────────────────────────────────
    # SLIDE 4 — I Wonder + Portfolio + Teamwork
    # ─────────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank)

    # Hero: full-bleed image with strong gradient
    has_img = _add_image_full_bleed(slide, sw, sh,
                                    f'{project.topic} india curiosity discovery',
                                    fallback_color=SLIDE_DEEP,
                                    variant=22 + challenge_n)
    if has_img:
        _add_dark_overlay(slide, Inches(0), Inches(0), sw, sh,
                          alpha=82, top_hex='1E1B4B', bottom_hex='000000')
    else:
        _add_rect(slide, Inches(0), Inches(0), sw, sh, fill=SLIDE_DEEP,
                  shape=MSO_SHAPE.RECTANGLE)

    # ── I Wonder block (top) ──
    _add_textbox(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.4),
                 '🔍  I WONDER...', font_size=14, bold=True, color=ACCENT_GOLD)
    _add_rect(slide, Inches(0.8), Inches(1.05), Inches(0.6),
              Inches(0.06), fill=ACCENT_GOLD, shape=MSO_SHAPE.RECTANGLE)
    _add_textbox(slide, Inches(0.8), Inches(1.25), Inches(11.7),
                 Inches(2.0), i_wonder or '', font_size=30, bold=True,
                 color=WHITE, italic=True, parse_md=True)

    # ── Portfolio task (middle) ──
    pf_card = _add_rect(slide, Inches(0.8), Inches(4.0), Inches(11.7),
                        Inches(1.6), fill=PRGBColor(0xFF, 0xFF, 0xFF),
                        shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _set_solid_fill_alpha(pf_card, 'FFFFFF', 14)
    tf = pf_card.text_frame
    tf.margin_left = Inches(0.45)
    tf.margin_top = Inches(0.25)
    tf.margin_right = Inches(0.45)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = '📓  FIND OUT MORE  ·  PORTFOLIO TASK'
    r.font.size = PPt(11)
    r.font.bold = True
    r.font.color.rgb = ACCENT_GOLD
    p2 = tf.add_paragraph()
    p2.space_before = PPt(6)
    _render_inline_pptx(p2, portfolio or '—', font_size=15,
                        color=WHITE, base_bold=True)

    # ── Teamwork bar (bottom) ──
    tw_card = _add_rect(slide, Inches(0.8), Inches(5.85), Inches(11.7),
                        Inches(1.0), fill=ACCENT_GOLD,
                        shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = tw_card.text_frame
    tf.margin_left = Inches(0.45)
    tf.margin_top = Inches(0.16)
    tf.margin_right = Inches(0.45)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = '🤝  WORK AS A TEAM'
    r.font.size = PPt(11)
    r.font.bold = True
    r.font.color.rgb = SLIDE_DEEP
    p2 = tf.add_paragraph()
    p2.space_before = PPt(4)
    r2 = p2.add_run()
    r2.text = ('Listen first  ·  speak up confidently  ·  '
               'disagree politely  ·  decide together')
    r2.font.size = PPt(13)
    r2.font.bold = True
    r2.font.color.rgb = SLIDE_DEEP

    # ── Footer ──
    _add_textbox(slide, Inches(0.8), sh - Inches(0.4), Inches(11.7),
                 Inches(0.35),
                 f'Challenge #{challenge_n}  ·  {grade}  ·  Neorise FSL',
                 font_size=10, color=PRGBColor(0xCB, 0xD5, 0xE1),
                 align=PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ╔══════════════════════════════════════════════════════════════════════╗
# ║  LESSON PLAN — built from scratch (DOCX)                             ║
# ╚══════════════════════════════════════════════════════════════════════╝

# Per-activity dot colors (5 hues that rotate). Matches the PPT theme palette
# but in docx RGBColor form.
_ACTIVITY_DOT_COLORS = [
    RGBColor(0xF5, 0x9E, 0x0B),  # Amber
    RGBColor(0x10, 0xB9, 0x81),  # Emerald
    RGBColor(0xEC, 0x48, 0x99),  # Rose
    RGBColor(0x4F, 0x46, 0xE5),  # Indigo
    RGBColor(0x06, 0xB6, 0xD4),  # Cyan
]

# Alternating row shade for tables
_ALT_ROW_SHADE = 'F8FAFC'


def _set_cell_padding(cell, *, top=0.15, bottom=0.15, left=0.2, right=0.2):
    """Set per-cell padding in cm via the tcMar element.
    1 cm == 567 twentieths-of-a-point (the dxa unit Word uses)."""
    tcPr = cell._tc.get_or_add_tcPr()
    tcMar = OxmlElement('w:tcMar')
    for side, val in [('top', top), ('left', left), ('bottom', bottom), ('right', right)]:
        node = OxmlElement(f'w:{side}')
        node.set(qn('w:w'), str(int(val * 567)))
        node.set(qn('w:type'), 'dxa')
        tcMar.append(node)
    tcPr.append(tcMar)


def _apply_table_padding(table, *, top=0.15, bottom=0.15, left=0.2, right=0.2):
    """Apply uniform cell padding to every cell in a table."""
    for row in table.rows:
        for cell in row.cells:
            _set_cell_padding(cell, top=top, bottom=bottom, left=left, right=right)


def _apply_alt_row_shading(table, *, header_rows=1, shade=_ALT_ROW_SHADE):
    """Apply a subtle background to odd-indexed body rows (skips header)."""
    for r_idx, row in enumerate(table.rows):
        if r_idx < header_rows:
            continue
        # Odd-index data rows get the shade (0-indexed AFTER header)
        if (r_idx - header_rows) % 2 == 1:
            for cell in row.cells:
                set_cell_shading(cell, shade)


def _cell_left_border(cell, color_hex='D97706', sz=24):
    """Thick left border on a single cell — used for callout boxes."""
    tcPr = cell._tc.get_or_add_tcPr()
    tcBorders = OxmlElement('w:tcBorders')
    left = OxmlElement('w:left')
    left.set(qn('w:val'), 'single')
    left.set(qn('w:sz'), str(sz))
    left.set(qn('w:color'), color_hex)
    tcBorders.append(left)
    # Hide the other borders so only left shows
    for side in ('top', 'right', 'bottom'):
        b = OxmlElement(f'w:{side}')
        b.set(qn('w:val'), 'nil')
        tcBorders.append(b)
    tcPr.append(tcBorders)


def _remove_table_borders(table):
    """Strip every border from a table (banner / accent line use this)."""
    tbl = table._element
    tblPr = tbl.find(qn('w:tblPr'))
    if tblPr is None:
        tblPr = OxmlElement('w:tblPr')
        tbl.insert(0, tblPr)
    # Remove any existing tblBorders
    existing = tblPr.find(qn('w:tblBorders'))
    if existing is not None:
        tblPr.remove(existing)
    tblBorders = OxmlElement('w:tblBorders')
    for side in ('top', 'left', 'bottom', 'right', 'insideH', 'insideV'):
        b = OxmlElement(f'w:{side}')
        b.set(qn('w:val'), 'nil')
        tblBorders.append(b)
    tblPr.append(tblBorders)


def _paragraph_top_border(paragraph, color_hex='E2E8F0', sz=6):
    """Thin horizontal rule above a paragraph (used before section headings)."""
    pPr = paragraph._p.get_or_add_pPr()
    pBdr = OxmlElement('w:pBdr')
    top = OxmlElement('w:top')
    top.set(qn('w:val'), 'single')
    top.set(qn('w:sz'), str(sz))
    top.set(qn('w:space'), '4')
    top.set(qn('w:color'), color_hex)
    pBdr.append(top)
    pPr.append(pBdr)


def _set_row_height_exact(table, row_idx, height_pt):
    """Lock a row to an exact height (Pt)."""
    tr = table.rows[row_idx]._tr
    trPr = tr.get_or_add_trPr()
    h = OxmlElement('w:trHeight')
    h.set(qn('w:val'), str(int(height_pt * 20)))
    h.set(qn('w:hRule'), 'exact')
    trPr.append(h)


def _add_cover_banner(doc, color_hex='630ED4', height_pt=18):
    """Full-width colored strip at the top of the cover page."""
    t = doc.add_table(rows=1, cols=1)
    t.autofit = False
    # Stretch to full page width (~16cm with 2.54cm margins on A4-ish page)
    try:
        t.columns[0].width = Cm(16)
    except Exception:
        pass
    cell = t.rows[0].cells[0]
    set_cell_shading(cell, color_hex)
    _set_cell_padding(cell, top=0, bottom=0, left=0, right=0)
    _set_row_height_exact(t, 0, height_pt)
    _remove_table_borders(t)
    # Clear the auto paragraph inside
    cell.paragraphs[0].text = ''


def _add_amber_accent_line(doc, color_hex='D97706', width_cm=4.0, height_pt=2):
    """Short thick amber line — used under the week title on cover."""
    t = doc.add_table(rows=1, cols=1)
    t.autofit = False
    try:
        t.columns[0].width = Cm(width_cm)
    except Exception:
        pass
    cell = t.rows[0].cells[0]
    set_cell_shading(cell, color_hex)
    _set_cell_padding(cell, top=0, bottom=0, left=0, right=0)
    _set_row_height_exact(t, 0, height_pt)
    _remove_table_borders(t)
    cell.paragraphs[0].text = ''


def _add_callout_box(doc, body_text, *, label='', bg_hex='FEF3C7',
                     border_hex='D97706', label_color=None,
                     body_color=None):
    """One-cell table styled as a colored callout (used for Facilitation
    Notes and Closure). Returns the cell for further customisation."""
    t = doc.add_table(rows=1, cols=1)
    t.autofit = True
    _remove_table_borders(t)
    cell = t.rows[0].cells[0]
    set_cell_shading(cell, bg_hex)
    _set_cell_padding(cell, top=0.18, bottom=0.18, left=0.35, right=0.25)
    _cell_left_border(cell, color_hex=border_hex, sz=28)
    p = cell.paragraphs[0]
    p.paragraph_format.space_after = Pt(0)
    if label:
        add_styled_run(p, label, bold=True, size=9,
                       color=label_color if label_color else MUTED)
        p2 = cell.add_paragraph()
        p2.paragraph_format.space_after = Pt(0)
        render_inline(p2, body_text, base_size=10,
                      base_color=body_color if body_color else INK)
    else:
        render_inline(p, body_text, base_size=10,
                      base_color=body_color if body_color else INK)
    # Small spacer after callout
    sp = doc.add_paragraph(); sp.paragraph_format.space_after = Pt(4)
    return cell


def _add_lesson_plan_header_footer(section, week: Week):
    """Right-aligned crumb in the header + centered page number in footer."""
    header_p = section.header.paragraphs[0]
    header_p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    # Clear any existing text
    for r in list(header_p.runs):
        r._r.getparent().remove(r._r)
    add_styled_run(header_p,
                   f'Neorise FSL · Week {week.number} · {week.phase}',
                   size=8, color=MUTED)

    footer_p = section.footer.paragraphs[0]
    footer_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for r in list(footer_p.runs):
        r._r.getparent().remove(r._r)
    run = footer_p.add_run()
    run.font.size = Pt(8)
    run.font.color.rgb = MUTED
    # Insert PAGE field
    fldBegin = OxmlElement('w:fldChar'); fldBegin.set(qn('w:fldCharType'), 'begin')
    instr = OxmlElement('w:instrText'); instr.text = ' PAGE '
    fldEnd = OxmlElement('w:fldChar'); fldEnd.set(qn('w:fldCharType'), 'end')
    run._r.append(fldBegin); run._r.append(instr); run._r.append(fldEnd)


def _render_teal_bullets(doc, text: str):
    """Render markdown bullet list with teal `•` markers instead of the
    default black List Bullet style. Non-bullet lines fall back to plain
    paragraphs."""
    if not text:
        return
    import re as _re
    lines = text.replace('\r\n', '\n').split('\n')
    for line in lines:
        s = line.strip()
        if not s:
            continue
        m = _re.match(r'^[\-\*•]\s+(.+)$', s)
        if m:
            p = doc.add_paragraph()
            p.paragraph_format.left_indent = Cm(0.6)
            p.paragraph_format.space_after = Pt(3)
            add_styled_run(p, '•  ', bold=True, size=11, color=ACCENT_TEAL)
            render_inline(p, m.group(1))
        else:
            # numbered list?
            mn = _re.match(r'^\d+\.\s+(.+)$', s)
            if mn:
                p = doc.add_paragraph()
                p.paragraph_format.left_indent = Cm(0.6)
                p.paragraph_format.space_after = Pt(3)
                add_styled_run(p, '•  ', bold=True, size=11, color=ACCENT_TEAL)
                render_inline(p, mn.group(1))
            else:
                p = doc.add_paragraph()
                p.paragraph_format.space_after = Pt(4)
                render_inline(p, s)


def build_lesson_plan_docx(project: Project, week: Week, content: dict) -> io.BytesIO:
    """Build the Weekly Lesson Plan DOCX in the client-approved layout.

    Sections per week:
      1. Cover (title + project meta table)
      2. Weekly Overview
      3. Week Challenge (3-row table: scenario / output / success criteria)
      4. Knowledge Focus (bullet list)
      5. Kaushal Bodh — Reflection Questions
      6. Competency Rubric (5-col table)
      7. Per block period:
         - Location row (5 cells with ☑ markers)
         - Materials Needed + Material Preparation
         - Glossary
         - Learning Objectives + Learning Outcomes (2-col table)
         - Activity Flow — for each activity:
            · Title + duration
            · Driving Focus + Expected Learning (2-col table)
            · Description
            · Discussion prompts
            · Facilitation notes
         - Closure
         - Portfolio Points
    """
    doc = Document()
    normal = doc.styles['Normal']
    normal.font.name = 'Calibri'
    normal.font.size = Pt(11)
    normal.font.color.rgb = INK
    normal.paragraph_format.line_spacing = 1.15  # match client reference

    for section in doc.sections:
        section.top_margin = Cm(2.54)   # 1 inch — matches client reference
        section.bottom_margin = Cm(2.54)
        section.left_margin = Cm(2.54)
        section.right_margin = Cm(2.54)
        # Page header + footer (#6)
        _add_lesson_plan_header_footer(section, week)

    # ── 1. Cover ──────────────────────────────────────────────────────
    title_para = doc.add_paragraph()
    title_para.paragraph_format.space_before = Pt(48)
    title_para.paragraph_format.space_after = Pt(4)
    add_styled_run(title_para, 'NEORISE FSL · WEEKLY LESSON PLAN',
                   bold=True, size=10, color=PRIMARY)

    p = doc.add_paragraph()
    p.paragraph_format.space_after = Pt(10)
    add_styled_run(p, project.topic, bold=True, size=26, color=INK)

    p = doc.add_paragraph()
    p.paragraph_format.space_after = Pt(6)
    add_styled_run(p, f'Week {week.number} — {week.phase}',
                   bold=True, size=15, color=ACCENT_AMBER)

    # (#1) Thin amber accent line under week title
    _add_amber_accent_line(doc, color_hex='D97706', width_cm=4.0, height_pt=2)

    sp = doc.add_paragraph()
    sp.paragraph_format.space_after = Pt(18)

    meta = doc.add_table(rows=0, cols=2)
    meta.autofit = True
    meta_rows = [
        ('GRADE',         project.grade),
        ('SUBJECT TRACK', project.get_subject_track_display()),
        ('TECH FOCUS',    ', '.join(project.tech_competency) or '—'),
        ('PHASE',         week.phase),
        ('BLOCK PERIODS', ', '.join(f'BP{s.number}' for s in week.sessions.order_by('number'))),
    ]
    for label, value in meta_rows:
        row = meta.add_row().cells
        row[0].width = Cm(4.5)
        set_cell_shading(row[0], 'F1F5F9')
        add_styled_run(row[0].paragraphs[0], label, bold=True, size=9, color=MUTED)
        add_styled_run(row[1].paragraphs[0], str(value), size=10, color=INK)
    _apply_table_padding(meta, top=0.15, bottom=0.15, left=0.2, right=0.2)

    doc.add_page_break()

    # ── 2. Weekly Overview ────────────────────────────────────────────
    _section_heading(doc, 'WEEKLY OVERVIEW')
    if content.get('weekly_overview'):
        render_markdown(doc, content['weekly_overview'])

    # ── 3. Week Challenge ─────────────────────────────────────────────
    challenge = content.get('week_challenge') or {}
    if challenge:
        _section_heading(doc, f'WEEK {week.number} CHALLENGE', space_before=14)
        ct = doc.add_table(rows=0, cols=2)
        ct.autofit = True
        rows = [
            ('1. Challenge',     challenge.get('scenario', '')),
            ('Output',           challenge.get('output', '')),
            ('Success Criteria', challenge.get('success_criteria', [])),
        ]
        for label, val in rows:
            r = ct.add_row().cells
            r[0].width = Cm(4.5)
            set_cell_shading(r[0], 'EDE9FE')
            add_styled_run(r[0].paragraphs[0], label,
                           bold=True, size=10, color=PRIMARY)
            if isinstance(val, list):
                first = True
                for v in val:
                    p = r[1].paragraphs[0] if first else r[1].add_paragraph()
                    first = False
                    add_styled_run(p, '• ', bold=True, size=10, color=ACCENT_TEAL)
                    render_inline(p, str(v))
            else:
                render_inline(r[1].paragraphs[0], str(val))
        # (#5) Padding on week challenge table
        _apply_table_padding(ct)

    # ── 4. Knowledge Focus ────────────────────────────────────────────
    if content.get('knowledge_focus'):
        _section_heading(doc, 'KNOWLEDGE FOCUS', space_before=14)
        # (#8) teal bullet markers instead of default black
        _render_teal_bullets(doc, content['knowledge_focus'])

    # ── 5. Kaushal Bodh — Reflection Questions ────────────────────────
    kb_questions = week.kaushal_bodh_questions or []
    if kb_questions:
        _section_heading(doc, 'KAUSHAL BODH — REFLECTION QUESTIONS', space_before=14)
        for q in kb_questions:
            p = doc.add_paragraph(style='List Bullet')
            p.paragraph_format.space_after = Pt(2)
            render_inline(p, str(q))

    # ── 6. Competency Rubric ──────────────────────────────────────────
    rubric = content.get('competency_rubric', [])
    if rubric:
        _section_heading(doc, 'COMPETENCY RUBRIC', space_before=14)
        rt = doc.add_table(rows=1, cols=5)
        rt.style = None
        rt.autofit = True
        head = rt.rows[0].cells
        for i, label in enumerate(['Code / Competency', 'Beginning', 'Developing', 'Proficient', 'Mastery']):
            set_cell_shading(head[i], 'EDE9FE')
            add_styled_run(head[i].paragraphs[0], label, bold=False, size=9, color=PRIMARY)
        for item in rubric:
            row = rt.add_row().cells
            cell_p = row[0].paragraphs[0]
            add_styled_run(cell_p, item.get('code', ''), bold=True, size=9, color=PRIMARY, font='Consolas')
            sub = row[0].add_paragraph()
            add_styled_run(sub, item.get('name', ''), size=9, color=INK)
            for j, level in enumerate(item.get('levels', [])[:4]):
                add_styled_run(row[j + 1].paragraphs[0], str(level), size=9, color=INK)
        # (#5) Padding + alt-row shading on rubric table
        _apply_table_padding(rt)
        _apply_alt_row_shading(rt, header_rows=1)

    # ── 7. Per-session breakdown ──────────────────────────────────────
    sessions_in_week = list(week.sessions.order_by('number'))
    session_data = content.get('sessions', [])

    for idx, session in enumerate(sessions_in_week):
        s_content = session_data[idx] if idx < len(session_data) else {}

        doc.add_page_break()

        # BP banner
        h = doc.add_paragraph()
        h.paragraph_format.space_after = Pt(0)
        add_styled_run(h, f'SESSION {session.number} · BLOCK PERIOD {session.number}',
                       bold=True, size=10, color=MUTED)

        title_text = s_content.get('title') or session.name
        h2 = doc.add_paragraph()
        h2.paragraph_format.space_after = Pt(8)
        add_styled_run(h2, title_text, bold=True, size=18, color=INK)

        if not s_content:
            warn = doc.add_paragraph()
            warn.paragraph_format.space_before = Pt(12)
            add_styled_run(
                warn,
                'Content generation for this block period failed. '
                'Click "Regenerate AI → Lesson Plan only" to retry.',
                italic=True, size=11, color=MUTED,
            )
            continue

        # Location checkbox row
        _location_row(doc, s_content.get('location', []))

        # Materials Needed
        materials = s_content.get('materials', [])
        if materials:
            _section_heading(doc, 'MATERIALS NEEDED', size=10, color=ACCENT_TEAL, space_before=10)
            for m in materials:
                p = doc.add_paragraph(style='List Bullet')
                p.paragraph_format.space_after = Pt(2)
                render_inline(p, str(m))

        # Learning Objectives + Outcomes (2-col table) — moved up to match reference
        objectives = s_content.get('learning_objectives', [])
        outcomes   = s_content.get('learning_outcomes', [])
        if objectives or outcomes:
            _section_heading(doc, 'LEARNING OBJECTIVES & OUTCOMES',
                             size=10, color=ACCENT_TEAL, space_before=10)
            t = doc.add_table(rows=1, cols=2)
            t.style = None
            t.autofit = True
            head = t.rows[0].cells
            set_cell_shading(head[0], 'EDE9FE')
            set_cell_shading(head[1], 'EDE9FE')
            add_styled_run(head[0].paragraphs[0], 'LEARNING OBJECTIVES',
                           bold=True, size=9, color=PRIMARY)
            add_styled_run(head[1].paragraphs[0], 'LEARNING OUTCOMES',
                           bold=True, size=9, color=PRIMARY)
            # One row per competency (matching client reference)
            max_rows = max(len(objectives), len(outcomes))
            for row_idx in range(max_rows):
                r = t.add_row().cells
                if row_idx < len(objectives):
                    obj = objectives[row_idx]
                    if isinstance(obj, dict):
                        sp_name  = obj.get('sp_name', '')
                        msp_code = obj.get('msp_code', '')
                        desc     = obj.get('description', '')
                    else:
                        sp_name, msp_code, desc = '', '', str(obj)
                    p = r[0].paragraphs[0]
                    if sp_name or msp_code:
                        add_styled_run(p, f'{sp_name} ', bold=True, size=9, color=PRIMARY)
                        if msp_code:
                            add_styled_run(p, f'| {msp_code} ', bold=True, size=9, color=ACCENT_AMBER, font='Consolas')
                    if desc:
                        add_styled_run(p, f'\n{desc}', size=9, color=INK)
                if row_idx < len(outcomes):
                    p = r[1].paragraphs[0]
                    add_styled_run(p, f'{row_idx + 1}. ', bold=True, size=9, color=ACCENT_TEAL)
                    render_inline(p, str(outcomes[row_idx]), base_size=9, base_color=INK)
            # (#5) Padding + alt-row shading on objectives table
            _apply_table_padding(t)
            _apply_alt_row_shading(t, header_rows=1)

        # Material Preparation
        prep = s_content.get('material_preparation', [])
        if prep:
            _section_heading(doc, 'MATERIAL PREPARATION', size=10, color=ACCENT_TEAL, space_before=10)
            for m in prep:
                p = doc.add_paragraph(style='List Bullet')
                p.paragraph_format.space_after = Pt(2)
                render_inline(p, str(m))

        # Glossary
        glossary = s_content.get('glossary', [])
        if glossary:
            _section_heading(doc, 'GLOSSARY', size=14, color=PRIMARY, space_before=8)
            for item in glossary:
                if isinstance(item, dict):
                    term = item.get('term', '')
                    defn = item.get('definition', '')
                else:
                    term, defn = '', str(item)
                p = doc.add_paragraph()
                p.paragraph_format.left_indent = Cm(0.4)
                p.paragraph_format.space_after = Pt(2)
                if term:
                    add_styled_run(p, f'{term}: ', bold=True, size=10, color=INK)
                render_inline(p, defn, base_size=10, base_color=INK)

        # Activity Flow
        activities = s_content.get('activities', [])
        if activities:
            _section_heading(doc, 'ACTIVITY FLOW', size=10, color=ACCENT_TEAL, space_before=12)
            for i, act in enumerate(activities, 1):
                _render_activity(doc, i, act)

        # (#4) Closure wrapped in light-purple quote box with purple left border
        if s_content.get('closure'):
            _section_heading(doc, 'CLOSURE', size=10, color=ACCENT_TEAL, space_before=10)
            _add_callout_box(
                doc,
                body_text=str(s_content['closure']),
                label='',                 # no label, just the quote body
                bg_hex='EDE9FE',          # light purple
                border_hex='630ED4',      # primary purple
                body_color=INK,
            )

        # Portfolio points — NOT rendered per session; consolidated at week end

    # ── 8. Weekly Portfolio Points (consolidated from both sessions) ───
    all_portfolio = []
    for idx, session in enumerate(sessions_in_week):
        s_content = session_data[idx] if idx < len(session_data) else {}
        pts = s_content.get('portfolio_points', [])
        if pts:
            all_portfolio.append((session, pts))

    if all_portfolio:
        doc.add_page_break()
        _section_heading(doc, 'WEEKLY PORTFOLIO POINTS', size=14, color=PRIMARY, space_before=12)
        p_intro = doc.add_paragraph()
        p_intro.paragraph_format.space_after = Pt(8)
        add_styled_run(p_intro,
                       'Record the following in your individual portfolio this week:',
                       italic=True, size=11, color=MUTED)

        for session, pts in all_portfolio:
            bp_label = doc.add_paragraph()
            bp_label.paragraph_format.space_before = Pt(8)
            bp_label.paragraph_format.space_after = Pt(4)
            add_styled_run(bp_label, f'Block Period {session.number}: {session.name}',
                           bold=True, size=11, color=ACCENT_TEAL)
            for pp in pts:
                p = doc.add_paragraph()
                p.paragraph_format.left_indent = Cm(0.6)
                p.paragraph_format.space_after = Pt(3)
                add_styled_run(p, '✎  ', bold=True, size=11, color=PRIMARY)
                render_inline(p, str(pp))

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


def _section_heading(doc, text, *, size=14, color=PRIMARY, space_before=8,
                     divider=True):
    """Section heading with an optional thin gray rule above (#2 design spec).
    The divider is suppressed for small subsection headings inside a BP
    (size <= 10) to avoid visual clutter."""
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(space_before)
    p.paragraph_format.space_after = Pt(4)
    if divider and size >= 12:
        _paragraph_top_border(p, color_hex='E2E8F0', sz=6)
        # Add some breathing room between rule and heading text
        p.paragraph_format.space_before = Pt(max(space_before, 12))
    add_styled_run(p, text, bold=True, size=size, color=color)


# Standard physical-location options for a session
_LOCATIONS = ['Lab', 'Activity Room', 'Outdoor', 'Classroom', 'Other Areas']


def _location_row(doc, selected: list):
    """5-col row showing which physical spaces this BP uses."""
    selected_lower = {str(s).strip().lower() for s in (selected or [])}
    t = doc.add_table(rows=1, cols=5)
    t.style = None
    t.autofit = True
    for i, loc in enumerate(_LOCATIONS):
        cell = t.rows[0].cells[i]
        on = loc.lower() in selected_lower
        if on:
            set_cell_shading(cell, 'D1FAE5')
        else:
            set_cell_shading(cell, 'F8FAFC')
        mark = '☑ ' if on else '☐ '
        add_styled_run(cell.paragraphs[0], mark + loc,
                       bold=on, size=9,
                       color=ACCENT_TEAL if on else MUTED)
    # (#5) Padding on location table
    _apply_table_padding(t, top=0.12, bottom=0.12, left=0.15, right=0.15)


def _render_activity(doc, idx: int, act: dict):
    """Render one activity block: title + driving-focus table + body."""
    # Visual spacer between activities (matches client reference)
    if idx > 1:
        spacer = doc.add_paragraph()
        spacer.style = doc.styles['Heading 2'] if 'Heading 2' in [s.name for s in doc.styles] else doc.styles['Normal']
        spacer.paragraph_format.space_before = Pt(6)
        spacer.paragraph_format.space_after = Pt(6)
        spacer.text = ''

    # Activity title — (#7) prefixed with a small colored dot
    dot_color = _ACTIVITY_DOT_COLORS[(idx - 1) % len(_ACTIVITY_DOT_COLORS)]
    ah = doc.add_paragraph()
    ah.paragraph_format.space_before = Pt(8)
    ah.paragraph_format.space_after = Pt(4)
    add_styled_run(ah, '● ', bold=True, size=14, color=dot_color)
    add_styled_run(ah, f'Activity {idx}: {act.get("name", "")}',
                   bold=True, size=12, color=PRIMARY)
    add_styled_run(ah, f'   ({act.get("duration", "—")})',
                   size=10, color=MUTED, italic=True)

    # Driving Focus + Expected Learning — 1-row 2-col with inline labels
    df = act.get('driving_focus', '')
    el = act.get('expected_learning', '')
    if df or el:
        t = doc.add_table(rows=1, cols=2)
        t.style = None
        t.autofit = True
        c = t.rows[0].cells
        set_cell_shading(c[0], 'FAF5FF')
        set_cell_shading(c[1], 'FAF5FF')
        p0 = c[0].paragraphs[0]
        add_styled_run(p0, 'Driving Focus : ', bold=True, size=10, color=PRIMARY)
        render_inline(p0, str(df), base_size=10, base_color=INK)
        p1 = c[1].paragraphs[0]
        add_styled_run(p1, 'Expected Learning : ', bold=True, size=10, color=PRIMARY)
        render_inline(p1, str(el), base_size=10, base_color=INK)
        # (#5) Cell padding
        _apply_table_padding(t, top=0.15, bottom=0.15, left=0.2, right=0.2)

    # Description
    if act.get('description'):
        render_markdown(doc, str(act['description']))

    # (#3) Facilitation notes wrapped in an amber callout box with left border
    if act.get('facilitation_notes'):
        _add_callout_box(
            doc,
            body_text=str(act['facilitation_notes']),
            label='FACILITATION NOTES',
            bg_hex='FEF3C7',          # light amber/cream
            border_hex='D97706',      # amber-600
            label_color=ACCENT_AMBER,
            body_color=INK,
        )


# ╔══════════════════════════════════════════════════════════════════════╗
# ║  SESSION PPT — modern, image-rich design                             ║
# ╚══════════════════════════════════════════════════════════════════════╝

# Modern brand palette (deep, contemporary tones)
SLIDE_PRIMARY = PRGBColor(0x4F, 0x46, 0xE5)   # Indigo 600
SLIDE_DEEP    = PRGBColor(0x1E, 0x1B, 0x4B)   # Indigo 950 — for hero overlays
SLIDE_INK     = PRGBColor(0x0F, 0x17, 0x2A)   # Slate 900
SLIDE_MUTED   = PRGBColor(0x64, 0x74, 0x8B)
SLIDE_BG      = PRGBColor(0xF8, 0xFA, 0xFC)   # Slate 50
SLIDE_PANEL   = PRGBColor(0xFF, 0xFF, 0xFF)
WHITE         = PRGBColor(0xFF, 0xFF, 0xFF)
SOFT_LINE     = PRGBColor(0xE2, 0xE8, 0xF0)
ACCENT_GOLD   = PRGBColor(0xFB, 0xBF, 0x24)

# Per-activity color theme — modern saturated tones with matching accents
ACTIVITY_THEMES = [
    {'main': PRGBColor(0xF5, 0x9E, 0x0B), 'tint': PRGBColor(0xFE, 0xF3, 0xC7),
     'deep': PRGBColor(0x78, 0x35, 0x0F), 'ink': PRGBColor(0x45, 0x20, 0x0B)},   # Amber
    {'main': PRGBColor(0x10, 0xB9, 0x81), 'tint': PRGBColor(0xD1, 0xFA, 0xE5),
     'deep': PRGBColor(0x06, 0x4E, 0x3B), 'ink': PRGBColor(0x06, 0x2F, 0x24)},   # Emerald
    {'main': PRGBColor(0xEC, 0x48, 0x99), 'tint': PRGBColor(0xFC, 0xE7, 0xF3),
     'deep': PRGBColor(0x83, 0x18, 0x43), 'ink': PRGBColor(0x50, 0x10, 0x29)},   # Rose
    {'main': PRGBColor(0x4F, 0x46, 0xE5), 'tint': PRGBColor(0xE0, 0xE7, 0xFF),
     'deep': PRGBColor(0x2E, 0x1B, 0x9C), 'ink': PRGBColor(0x1E, 0x1B, 0x4B)},   # Indigo
    {'main': PRGBColor(0x06, 0xB6, 0xD4), 'tint': PRGBColor(0xCF, 0xFA, 0xFE),
     'deep': PRGBColor(0x0E, 0x49, 0x66), 'ink': PRGBColor(0x08, 0x36, 0x44)},   # Cyan
]


# ── XML helpers for transparency / gradients ──────────────────────────────
def _set_solid_fill_alpha(shape, rgb_hex: str, alpha_pct: int):
    """Solid fill with alpha (0-100). Manipulates underlying XML."""
    sp = shape.fill._xPr
    # Remove existing fill children
    for child in list(sp):
        tag = child.tag.split('}')[-1]
        if tag.endswith('Fill'):
            sp.remove(child)
    fill_xml = (
        f'<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:srgbClr val="{rgb_hex}">'
        f'<a:alpha val="{int(alpha_pct * 1000)}"/>'
        f'</a:srgbClr>'
        f'</a:solidFill>'
    )
    sp.insert(0, parse_xml(fill_xml))


def _set_gradient_fill(shape, top_hex: str, bottom_hex: str,
                       top_alpha: int = 100, bottom_alpha: int = 100):
    """Linear vertical gradient (top → bottom). Alphas in percent."""
    sp = shape.fill._xPr
    for child in list(sp):
        tag = child.tag.split('}')[-1]
        if tag.endswith('Fill'):
            sp.remove(child)
    grad_xml = f'''<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
        flip="none" rotWithShape="1">
      <a:gsLst>
        <a:gs pos="0"><a:srgbClr val="{top_hex}"><a:alpha val="{top_alpha*1000}"/></a:srgbClr></a:gs>
        <a:gs pos="100000"><a:srgbClr val="{bottom_hex}"><a:alpha val="{bottom_alpha*1000}"/></a:srgbClr></a:gs>
      </a:gsLst>
      <a:lin ang="5400000" scaled="0"/>
    </a:gradFill>'''
    sp.insert(0, parse_xml(grad_xml))


# ── Inline markdown for pptx (bold + italic) ──────────────────────────────
_PPTX_INLINE_RE = re.compile(r'(\*\*[^*\n]+?\*\*|\*[^*\n]+?\*)')


def _render_inline_pptx(paragraph, text, *, font_size=14, color=SLIDE_INK,
                        base_bold=False, italic=False):
    """Render text into a pptx paragraph, splitting bold/italic spans into runs."""
    if text is None:
        text = ''
    # Clear any existing runs in this paragraph
    for run in list(paragraph.runs):
        run._r.getparent().remove(run._r)

    parts = _PPTX_INLINE_RE.split(str(text))
    for part in parts:
        if not part:
            continue
        run = paragraph.add_run()
        if part.startswith('**') and part.endswith('**'):
            run.text = part[2:-2]
            run.font.bold = True
            run.font.italic = italic
        elif part.startswith('*') and part.endswith('*'):
            run.text = part[1:-1]
            run.font.bold = base_bold
            run.font.italic = True
        else:
            run.text = part
            run.font.bold = base_bold
            run.font.italic = italic
        run.font.size = PPt(font_size)
        run.font.color.rgb = color


def _add_textbox(slide, x, y, w, h, text, *, font_size=18, bold=False,
                 color=SLIDE_INK, align=None, italic=False, parse_md=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    if parse_md:
        _render_inline_pptx(p, text, font_size=font_size, color=color,
                            base_bold=bold, italic=italic)
    else:
        run = p.add_run()
        run.text = str(text)
        run.font.size = PPt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return box


def _add_rect(slide, x, y, w, h, *, fill=None, line=None,
              line_width_pt=0.75, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    rect = slide.shapes.add_shape(shape, x, y, w, h)
    rect.shadow.inherit = False
    if fill is not None:
        rect.fill.solid()
        rect.fill.fore_color.rgb = fill
    else:
        rect.fill.background()
    if line is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line
        rect.line.width = PPt(line_width_pt)
    rect.text_frame.text = ''
    return rect


def _add_title_bar(slide, slide_w, color, label, *, sublabel=None):
    """Coloured title bar across the top with a label and optional sub-label."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, Inches(0.7))
    bar.shadow.inherit = False
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = label
    run.font.size = PPt(20)
    run.font.bold = True
    run.font.color.rgb = WHITE
    if sublabel:
        sp = tf.add_paragraph()
        srun = sp.add_run()
        srun.text = sublabel
        srun.font.size = PPt(11)
        srun.font.bold = True
        srun.font.color.rgb = PRGBColor(0xFF, 0xE4, 0xC0)


def _set_slide_background(slide, color=SLIDE_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _new_slide(prs, blank_layout, bg_color=SLIDE_BG):
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_background(slide, bg_color)
    return slide


# ── Query helpers ─────────────────────────────────────────────────────────
def _activity_image_query(act: dict, fallback: str = '') -> str:
    """Best-quality search query for an activity slide. Priority:
       1. Explicit `image_query` field (AI-provided, optimised for search)
       2. `media_placeholder` (AI-described visual)
       3. Activity title only (no project topic — too generic in search)
       4. Fallback string"""
    if act.get('image_query'):
        return str(act['image_query'])
    if act.get('media_placeholder'):
        return str(act['media_placeholder'])
    if act.get('title'):
        return str(act['title'])
    return fallback


def _session_image_query(content: dict, project, session, kind: str) -> str:
    """Top-level cover/closing query. Falls back to project topic."""
    explicit = content.get(f'{kind}_image_query') or content.get('image_query')
    if explicit:
        return str(explicit)
    if kind == 'closing':
        return f'inspiring landscape sunrise {project.topic}'
    # cover
    return f'{content.get("title") or session.name} {project.topic}'


# ── Caption cleaner ───────────────────────────────────────────────────────
def _clean_caption(text: str, max_chars: int = 110) -> str:
    """Strip prompt-like meta phrases and return a short, readable caption."""
    if not text:
        return ''
    s = str(text).strip()
    # Drop "Source: ..." / "Search ... on ..." trailing fragments
    s = re.sub(r'\s*\bsource\s*:.*$', '', s, flags=re.IGNORECASE)
    s = re.sub(r'\s*\bsearch\s+["\'].+$', '', s, flags=re.IGNORECASE)
    s = re.sub(r'@\w+', '', s)
    s = re.sub(r'https?://\S+', '', s)
    # Strip leading verb prefixes like "Display ...", "Show ...", "Use ..."
    s = re.sub(r'^(?:display|show|use|find|source|search|suggested|please)\s+'
               r'(?:a\s+)?(?:high[\s\-]?(?:resolution|res|quality|definition)\s+)?'
               r'(?:photo|image|picture|video|clip|footage|graphic|illustration|'
               r'infographic|montage|collage|reel|short)s?[\s:\-—]*',
               '', s, flags=re.IGNORECASE)
    # Strip leading noun prefixes: "Photo of ...", "Short video (60s):", etc.
    s = re.sub(r'^(?:a\s+)?(?:short|long|brief|quick|high[\s\-]?(?:resolution|res|quality|definition))\s+'
               r'(?:photo|image|picture|video|clip|footage|graphic|illustration|'
               r'infographic|montage|collage|reel)s?\b',
               '', s, flags=re.IGNORECASE).lstrip(' :-—')
    s = re.sub(r'^(?:photo|image|picture|video|clip|footage|graphic|illustration|'
               r'infographic|montage|collage)s?\s+(?:of|showing|featuring|depicting)\s+'
               r'(?:a\s+|an\s+|the\s+)?',
               '', s, flags=re.IGNORECASE)
    # Catch leftover "of/showing/featuring" at the start (after a noun was stripped)
    s = re.sub(r'^(?:of|showing|featuring|depicting)\s+(?:a\s+|an\s+|the\s+)?',
               '', s, flags=re.IGNORECASE)
    # Strip trailing parenthetical durations "(60-90 sec)"
    s = re.sub(r'\(\s*\d+[\s\-–]*\d*\s*(?:sec|second|min|minute)s?\.?\s*\)\s*[:\-—]?\s*',
               '', s, flags=re.IGNORECASE)
    s = re.sub(r'\s+', ' ', s).strip(' .,—-:;')
    if s and s[0].islower():
        s = s[0].upper() + s[1:]
    if len(s) > max_chars:
        s = s[:max_chars].rsplit(' ', 1)[0].rstrip(' ,;:-—') + '…'
    return s


def _crop_to_aspect(img_io, target_w_emu, target_h_emu):
    """Center-crop a BytesIO image to match a target aspect ratio so it fills
    the slide box without stretching. Returns a fresh BytesIO."""
    if not _PIL_OK or img_io is None:
        return img_io
    try:
        img_io.seek(0)
        im = _PILImage.open(img_io)
        # Convert palette / RGBA → RGB for clean JPEG saving
        if im.mode != 'RGB':
            im = im.convert('RGB')
        src_w, src_h = im.size
        target_aspect = target_w_emu / max(target_h_emu, 1)
        src_aspect = src_w / src_h
        if abs(src_aspect - target_aspect) < 0.02:
            # Already close enough — skip crop
            img_io.seek(0)
            return img_io
        if src_aspect > target_aspect:
            # Source is wider → crop sides
            new_w = int(src_h * target_aspect)
            x = (src_w - new_w) // 2
            box = (x, 0, x + new_w, src_h)
        else:
            # Source is taller → crop top/bottom (favouring upper third for faces)
            new_h = int(src_w / target_aspect)
            y = max(0, (src_h - new_h) // 2 - int(new_h * 0.05))
            box = (0, y, src_w, y + new_h)
        cropped = im.crop(box)
        out = io.BytesIO()
        cropped.save(out, format='JPEG', quality=88, optimize=True)
        out.seek(0)
        return out
    except Exception:
        img_io.seek(0)
        return img_io


# ── Image helpers ─────────────────────────────────────────────────────────
def _add_image_full_bleed(slide, sw, sh, query: str, *, fallback_color=None,
                          variant: int = 0):
    """Fill the whole slide with a stock image (cover-fit). Returns True if
    image was added, False on fallback."""
    img = fetch_stock_image_io(query, orientation='landscape',
                               variant=variant) if query else None
    if img is None:
        if fallback_color is not None:
            _add_rect(slide, Inches(0), Inches(0), sw, sh,
                      fill=fallback_color, shape=MSO_SHAPE.RECTANGLE)
        return False
    try:
        cropped = _crop_to_aspect(img, sw, sh)
        slide.shapes.add_picture(cropped, Inches(0), Inches(0),
                                 width=sw, height=sh)
        return True
    except Exception:
        if fallback_color is not None:
            _add_rect(slide, Inches(0), Inches(0), sw, sh,
                      fill=fallback_color, shape=MSO_SHAPE.RECTANGLE)
        return False


def _add_image_box(slide, x, y, w, h, query: str, *, fallback_color=None,
                   variant: int = 0):
    """Embed a stock image in a defined box. Returns True if image added."""
    img = fetch_stock_image_io(query, variant=variant) if query else None
    if img is None:
        if fallback_color is not None:
            _add_rect(slide, x, y, w, h, fill=fallback_color,
                      shape=MSO_SHAPE.RECTANGLE)
        return False
    try:
        cropped = _crop_to_aspect(img, w, h)
        slide.shapes.add_picture(cropped, x, y, width=w, height=h)
        return True
    except Exception:
        if fallback_color is not None:
            _add_rect(slide, x, y, w, h, fill=fallback_color,
                      shape=MSO_SHAPE.RECTANGLE)
        return False


def _add_dark_overlay(slide, x, y, w, h, *, alpha=55, top_hex='000000',
                      bottom_hex='000000'):
    """Semi-transparent dark gradient overlay for image legibility."""
    rect = _add_rect(slide, x, y, w, h, fill=PRGBColor(0, 0, 0),
                     shape=MSO_SHAPE.RECTANGLE)
    _set_gradient_fill(rect, top_hex, bottom_hex,
                       top_alpha=max(0, alpha - 25), bottom_alpha=alpha + 20)
    return rect


def _add_decorative_dot(slide, x, y, d, color):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    dot.shadow.inherit = False
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()


# ── Per-slide builders (modern design) ────────────────────────────────────
def _slide_cover(prs, layout, sw, sh, project, session, content):
    slide = prs.slides.add_slide(layout)
    hero_query = _session_image_query(content, project, session, 'cover')
    # Variant 0 reserved for cover so it stays distinctive
    has_img = _add_image_full_bleed(slide, sw, sh, hero_query,
                                    fallback_color=SLIDE_DEEP, variant=0)
    # Dark gradient overlay for legibility
    if has_img:
        _add_dark_overlay(slide, Inches(0), Inches(0), sw, sh,
                          alpha=70, top_hex='1E1B4B', bottom_hex='000000')
    else:
        _add_rect(slide, Inches(0), Inches(0), sw, sh,
                  fill=SLIDE_DEEP, shape=MSO_SHAPE.RECTANGLE)

    # Top eyebrow
    _add_textbox(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.4),
                 'NEORISE FSL  ·  STUDENT DECK', font_size=11, bold=True,
                 color=ACCENT_GOLD)

    # Number tag (top right)
    tag = _add_rect(slide, Inches(11.4), Inches(0.55), Inches(1.4),
                    Inches(0.55), fill=ACCENT_GOLD,
                    shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = tag.text_frame
    tf.margin_left = Inches(0)
    tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f'BP {session.number}'
    run.font.size = PPt(16)
    run.font.bold = True
    run.font.color.rgb = SLIDE_DEEP

    # Bottom-left content stack
    _add_textbox(slide, Inches(0.8), Inches(4.4), Inches(11.5), Inches(0.5),
                 f'WEEK {session.week.number}  ·  {session.week.phase.upper()}',
                 font_size=14, bold=True, color=ACCENT_GOLD)
    _add_textbox(slide, Inches(0.8), Inches(4.95), Inches(11.5), Inches(1.6),
                 content.get('title') or session.name,
                 font_size=44, bold=True, color=WHITE)
    _add_textbox(slide, Inches(0.8), Inches(6.55), Inches(11.5), Inches(0.45),
                 project.topic, font_size=16, color=PRGBColor(0xCB, 0xD5, 0xE1))
    _add_textbox(slide, Inches(0.8), Inches(7.0), Inches(11.5), Inches(0.4),
                 f'Grade {project.grade}  ·  {project.get_subject_track_display()}',
                 font_size=12, color=PRGBColor(0x94, 0xA3, 0xB8))


def _slide_goals(prs, layout, sw, sh, content):
    goals = content.get('goals', [])
    if not goals:
        return
    slide = _new_slide(prs, layout)

    # Eyebrow + headline
    _add_textbox(slide, Inches(0.8), Inches(0.6), Inches(12), Inches(0.4),
                 'TODAY · OUR GOALS', font_size=11, bold=True, color=SLIDE_PRIMARY)
    _add_textbox(slide, Inches(0.8), Inches(1.0), Inches(12), Inches(0.9),
                 'What you will walk away with', font_size=34, bold=True,
                 color=SLIDE_INK)
    # Accent rule
    _add_rect(slide, Inches(0.8), Inches(2.05), Inches(0.8), Inches(0.08),
              fill=ACCENT_GOLD, shape=MSO_SHAPE.RECTANGLE)

    gy = Inches(2.55)
    for i, goal in enumerate(goals[:5]):
        # Big number
        _add_textbox(slide, Inches(0.8), gy, Inches(1.4), Inches(0.85),
                     f'0{i+1}' if i < 9 else str(i+1), font_size=48,
                     bold=True, color=PRGBColor(0xCB, 0xD5, 0xE1))
        # Goal text
        gbox = _add_textbox(slide, Inches(2.4), gy + Inches(0.18),
                            Inches(10.3), Inches(0.7), '',
                            font_size=20, bold=False, color=SLIDE_INK,
                            parse_md=True)
        _render_inline_pptx(gbox.text_frame.paragraphs[0], goal,
                            font_size=20, color=SLIDE_INK, base_bold=True)
        # Thin separator
        if i < min(len(goals), 5) - 1:
            _add_rect(slide, Inches(0.8), gy + Inches(0.95), Inches(12),
                      Inches(0.015), fill=SOFT_LINE, shape=MSO_SHAPE.RECTANGLE)
        gy += Inches(1.0)


def _slide_timeline(prs, layout, sw, sh, timeline):
    if not timeline:
        return
    slide = _new_slide(prs, layout)

    _add_textbox(slide, Inches(0.8), Inches(0.6), Inches(12), Inches(0.4),
                 '⏱  THE FLOW', font_size=11, bold=True, color=SLIDE_PRIMARY)
    _add_textbox(slide, Inches(0.8), Inches(1.0), Inches(12), Inches(0.9),
                 "How our 80 minutes will flow", font_size=32, bold=True,
                 color=SLIDE_INK)
    _add_rect(slide, Inches(0.8), Inches(2.0), Inches(0.8), Inches(0.08),
              fill=ACCENT_GOLD, shape=MSO_SHAPE.RECTANGLE)

    # Vertical timeline bar
    _add_rect(slide, Inches(2.05), Inches(2.5), Inches(0.06), Inches(4.3),
              fill=PRGBColor(0xCB, 0xD5, 0xE1), shape=MSO_SHAPE.RECTANGLE)

    ty = Inches(2.4)
    for i, item in enumerate(timeline[:6]):
        theme = ACTIVITY_THEMES[i % len(ACTIVITY_THEMES)]
        # Dot on the bar
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.85), ty,
                                     Inches(0.5), Inches(0.5))
        dot.shadow.inherit = False
        dot.fill.solid()
        dot.fill.fore_color.rgb = theme['main']
        dot.line.color.rgb = WHITE
        dot.line.width = PPt(2)
        # Time
        _add_textbox(slide, Inches(0.8), ty - Inches(0.05), Inches(1.0),
                     Inches(0.5), str(item.get('time', '')), font_size=12,
                     bold=True, color=theme['ink'])
        # Activity card
        card = _add_rect(slide, Inches(2.6), ty - Inches(0.1), Inches(10.2),
                         Inches(0.7), fill=WHITE, line=SOFT_LINE,
                         line_width_pt=0.75)
        tf = card.text_frame
        tf.margin_left = Inches(0.3)
        tf.margin_top = Inches(0.18)
        tf.word_wrap = True
        _render_inline_pptx(tf.paragraphs[0], str(item.get('activity', '')),
                            font_size=15, color=SLIDE_INK, base_bold=True)
        ty += Inches(0.75)


def _slide_activity_intro(prs, layout, sw, sh, idx, act, theme, project):
    """Modern hero-style intro: full-bleed image + dark overlay + giant number."""
    slide = prs.slides.add_slide(layout)
    query = _activity_image_query(act, fallback=project.topic)
    # Unique variant per activity intro: 1, 4, 7, 10... (gap of 3 from siblings)
    has_img = _add_image_full_bleed(slide, sw, sh, query,
                                    fallback_color=theme['deep'],
                                    variant=1 + idx * 3)
    # Dark gradient overlay
    if has_img:
        deep_hex = '{:02X}{:02X}{:02X}'.format(theme['deep'][0],
                                                theme['deep'][1],
                                                theme['deep'][2])
        _add_dark_overlay(slide, Inches(0), Inches(0), sw, sh,
                          alpha=78, top_hex=deep_hex, bottom_hex='000000')
    else:
        _add_rect(slide, Inches(0), Inches(0), sw, sh,
                  fill=theme['deep'], shape=MSO_SHAPE.RECTANGLE)

    # Decorative gold accent line
    _add_rect(slide, Inches(0.8), Inches(2.1), Inches(0.8), Inches(0.08),
              fill=ACCENT_GOLD, shape=MSO_SHAPE.RECTANGLE)

    # Eyebrow
    _add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
                 f'ACTIVITY {idx + 1}  OF  5', font_size=14, bold=True,
                 color=ACCENT_GOLD)

    # Big title
    _add_textbox(slide, Inches(0.8), Inches(2.4), Inches(11.5), Inches(2.6),
                 act.get('title', 'Activity'), font_size=56, bold=True,
                 color=WHITE, parse_md=True)

    # Duration pill
    duration = act.get('duration')
    if duration:
        pill = _add_rect(slide, Inches(0.8), Inches(5.2), Inches(2.6),
                         Inches(0.7), fill=ACCENT_GOLD,
                         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tf = pill.text_frame
        tf.margin_left = Inches(0)
        tf.margin_top = Inches(0.12)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f'⏱  {duration}'
        run.font.size = PPt(16)
        run.font.bold = True
        run.font.color.rgb = SLIDE_DEEP

    # Giant outlined number bottom-right
    _add_textbox(slide, Inches(10.5), Inches(4.8), Inches(2.5), Inches(2.4),
                 str(idx + 1), font_size=240, bold=True,
                 color=PRGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.RIGHT)


def _slide_activity_content(prs, layout, sw, sh, idx, act, theme, blocks,
                            project):
    """Modern split layout: image strip on left + content cards right.
    First content slide carries an image; subsequent ones are text-only."""
    if not blocks:
        return
    slide = _new_slide(prs, layout)

    # Eyebrow
    _add_textbox(slide, Inches(0.8), Inches(0.5), Inches(12), Inches(0.4),
                 f'ACTIVITY {idx + 1}', font_size=11, bold=True,
                 color=theme['main'])
    # Headline
    _add_textbox(slide, Inches(0.8), Inches(0.9), Inches(12), Inches(0.7),
                 act.get('title', ''), font_size=26, bold=True,
                 color=SLIDE_INK, parse_md=True)
    _add_rect(slide, Inches(0.8), Inches(1.7), Inches(0.6), Inches(0.06),
              fill=theme['main'], shape=MSO_SHAPE.RECTANGLE)

    # Image on the left — different variant from intro to avoid duplicate
    img_query = _activity_image_query(act, fallback=project.topic)
    img_added = _add_image_box(slide, Inches(0.8), Inches(2.1), Inches(4.5),
                               Inches(4.7), img_query,
                               fallback_color=theme['tint'],
                               variant=2 + idx * 3)

    content_x = Inches(5.6) if img_added else Inches(0.8)
    content_w = Inches(7.2) if img_added else Inches(12)

    y = Inches(2.1)
    for cb in blocks:
        # Accent dot
        _add_decorative_dot(slide, content_x, y + Inches(0.18),
                            Inches(0.18), theme['main'])
        # Text
        box = _add_textbox(slide, content_x + Inches(0.35), y, content_w,
                           Inches(1.4), '', font_size=15, color=SLIDE_INK)
        tf = box.text_frame
        tf.word_wrap = True
        _render_inline_pptx(tf.paragraphs[0], cb, font_size=15,
                            color=SLIDE_INK)
        y += Inches(1.45)


def _slide_activity_prompts(prs, layout, sw, sh, idx, act, theme, prompts):
    """Modern prompts slide — clean numbered list with subtle theme accent."""
    if not prompts:
        return
    slide = _new_slide(prs, layout)

    # Eyebrow + headline
    _add_textbox(slide, Inches(0.8), Inches(0.6), Inches(12), Inches(0.4),
                 f'ACTIVITY {idx + 1}  ·  DISCUSS', font_size=11, bold=True,
                 color=theme['main'])
    _add_textbox(slide, Inches(0.8), Inches(1.0), Inches(12), Inches(0.9),
                 'Talk it out with your group', font_size=32, bold=True,
                 color=SLIDE_INK)
    _add_rect(slide, Inches(0.8), Inches(2.0), Inches(0.8), Inches(0.08),
              fill=theme['main'], shape=MSO_SHAPE.RECTANGLE)

    y = Inches(2.55)
    for i, pr in enumerate(prompts[:5]):
        # Big number
        _add_textbox(slide, Inches(0.8), y - Inches(0.1), Inches(1.0),
                     Inches(0.7), f'{i+1:02d}', font_size=36, bold=True,
                     color=theme['main'])
        # Prompt text
        box = _add_textbox(slide, Inches(2.0), y, Inches(10.7), Inches(0.7),
                           '', font_size=18, color=SLIDE_INK, parse_md=True)
        _render_inline_pptx(box.text_frame.paragraphs[0], pr, font_size=18,
                            color=SLIDE_INK, base_bold=True)
        # Separator
        if i < min(len(prompts), 5) - 1:
            _add_rect(slide, Inches(2.0), y + Inches(0.85),
                      Inches(10.7), Inches(0.015), fill=SOFT_LINE,
                      shape=MSO_SHAPE.RECTANGLE)
        y += Inches(0.95)


def _slide_media_image(prs, layout, sw, sh, idx, act, theme, media_text):
    """Hero image slide using the AI's media suggestion."""
    slide = prs.slides.add_slide(layout)

    # Use cleaned caption as the actual image search query (drops "Display ..." etc.)
    search_query = _clean_caption(media_text, max_chars=200) or media_text
    has_img = _add_image_full_bleed(slide, sw, sh, search_query,
                                    fallback_color=theme['tint'],
                                    variant=3 + idx * 3)

    caption = _clean_caption(media_text, max_chars=110)

    if has_img:
        # Stronger gradient at bottom for caption legibility
        _add_dark_overlay(slide, Inches(0), Inches(5.0), sw, Inches(2.5),
                          alpha=88, top_hex='000000', bottom_hex='000000')
        # Eyebrow
        _add_textbox(slide, Inches(0.8), Inches(5.6), Inches(11.5),
                     Inches(0.4), f'ACTIVITY {idx + 1}  ·  LOOK CLOSELY',
                     font_size=11, bold=True, color=ACCENT_GOLD)
        # Caption — single line, auto-shrunk
        _add_textbox(slide, Inches(0.8), Inches(6.05), Inches(11.5),
                     Inches(1.1), caption, font_size=18, bold=True,
                     color=WHITE, parse_md=True)
    else:
        # Fallback placeholder card
        _set_slide_background(slide, theme['tint'])
        _add_textbox(slide, Inches(0.8), Inches(0.6), Inches(12),
                     Inches(0.4), f'ACTIVITY {idx + 1}  ·  VISUAL',
                     font_size=11, bold=True, color=theme['main'])
        _add_textbox(slide, Inches(0.8), Inches(1.0), Inches(12),
                     Inches(0.9), 'Show this to the class', font_size=30,
                     bold=True, color=SLIDE_INK)
        _add_rect(slide, Inches(0.8), Inches(2.0), Inches(0.8),
                  Inches(0.08), fill=theme['main'],
                  shape=MSO_SHAPE.RECTANGLE)
        ph = _add_rect(slide, Inches(1.5), Inches(2.6), Inches(10.3),
                       Inches(3.6), fill=WHITE, line=theme['main'],
                       line_width_pt=2.5)
        tf = ph.text_frame
        tf.margin_top = Inches(1.2)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = '🖼️  ADD IMAGE HERE'
        r.font.size = PPt(28)
        r.font.bold = True
        r.font.color.rgb = theme['main']
        _add_textbox(slide, Inches(1.5), Inches(6.4), Inches(10.3),
                     Inches(0.7), caption, font_size=14, italic=True,
                     color=SLIDE_MUTED, align=PP_ALIGN.CENTER, parse_md=True)


def _slide_video_placeholder(prs, layout, sw, sh, idx, act, theme,
                             video_text, project):
    """Video player look-alike slide. Uses an image background with a big
    play-button overlay + caption — mimics a real video embed."""
    slide = prs.slides.add_slide(layout)

    search_query = _clean_caption(video_text, max_chars=200) or video_text
    has_img = _add_image_full_bleed(slide, sw, sh, search_query,
                                    fallback_color=PRGBColor(0x12, 0x12, 0x16),
                                    variant=4 + idx * 3)

    # Strong dark overlay so play button + caption pop
    _add_dark_overlay(slide, Inches(0), Inches(0), sw, sh,
                      alpha=72, top_hex='0F172A', bottom_hex='000000')

    # Top-left "VIDEO" tag
    tag = _add_rect(slide, Inches(0.8), Inches(0.6), Inches(1.7),
                    Inches(0.55), fill=PRGBColor(0xDC, 0x26, 0x26),
                    shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = tag.text_frame
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = '● VIDEO'
    r.font.size = PPt(14)
    r.font.bold = True
    r.font.color.rgb = WHITE

    # Eyebrow
    _add_textbox(slide, Inches(0.8), Inches(1.4), Inches(11.5),
                 Inches(0.4), f'ACTIVITY {idx + 1}  ·  WATCH',
                 font_size=12, bold=True, color=ACCENT_GOLD)

    # Big play button (white circle + dark triangle)
    cx = sw / 2 - Inches(0.9)
    cy = sh / 2 - Inches(1.4)
    play_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy,
                                         Inches(1.8), Inches(1.8))
    play_circle.shadow.inherit = False
    play_circle.fill.solid()
    play_circle.fill.fore_color.rgb = PRGBColor(0xFF, 0xFF, 0xFF)
    play_circle.line.fill.background()

    triangle = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,
                                      cx + Inches(0.55), cy + Inches(0.45),
                                      Inches(0.85), Inches(0.9))
    triangle.shadow.inherit = False
    triangle.fill.solid()
    triangle.fill.fore_color.rgb = PRGBColor(0x0F, 0x17, 0x2A)
    triangle.line.fill.background()
    triangle.rotation = 30

    # "Click to play" hint under the button
    _add_textbox(slide, Inches(0.5), cy + Inches(2.1), sw - Inches(1),
                 Inches(0.4), '▶  Play in class — placeholder for teacher',
                 font_size=12, color=PRGBColor(0xCB, 0xD5, 0xE1),
                 align=PP_ALIGN.CENTER, italic=True)

    # Bottom caption strip
    _add_dark_overlay(slide, Inches(0), Inches(6.0), sw, Inches(1.5),
                      alpha=92, top_hex='000000', bottom_hex='000000')
    caption = _clean_caption(video_text, max_chars=120)
    _add_textbox(slide, Inches(0.8), Inches(6.45), Inches(11.5),
                 Inches(0.95), caption, font_size=18, bold=True,
                 color=WHITE, parse_md=True)


def _slide_reflection(prs, layout, sw, sh, reflections):
    if not reflections:
        return
    slide = _new_slide(prs, layout)

    _add_textbox(slide, Inches(0.8), Inches(0.6), Inches(12), Inches(0.4),
                 'REFLECT  ·  WRAP IT UP', font_size=11, bold=True,
                 color=SLIDE_PRIMARY)
    _add_textbox(slide, Inches(0.8), Inches(1.0), Inches(12), Inches(0.9),
                 'Take a minute. Write it down.', font_size=32, bold=True,
                 color=SLIDE_INK)
    _add_rect(slide, Inches(0.8), Inches(2.0), Inches(0.8), Inches(0.08),
              fill=ACCENT_GOLD, shape=MSO_SHAPE.RECTANGLE)

    y = Inches(2.55)
    for i, r in enumerate(reflections[:5]):
        # Pencil-style accent
        _add_decorative_dot(slide, Inches(0.85), y + Inches(0.22),
                            Inches(0.18), SLIDE_PRIMARY)
        box = _add_textbox(slide, Inches(1.3), y, Inches(11.5),
                           Inches(0.65), '', font_size=19, color=SLIDE_INK,
                           parse_md=True)
        _render_inline_pptx(box.text_frame.paragraphs[0], r, font_size=19,
                            color=SLIDE_INK, base_bold=True)
        if i < min(len(reflections), 5) - 1:
            _add_rect(slide, Inches(1.3), y + Inches(0.78),
                      Inches(11.5), Inches(0.015), fill=SOFT_LINE,
                      shape=MSO_SHAPE.RECTANGLE)
        y += Inches(0.85)


def _slide_closing(prs, layout, sw, sh, content, project, session):
    closing = content.get('closing_thought') if isinstance(content, dict) else content
    if not closing:
        return
    slide = prs.slides.add_slide(layout)

    # Hero background — prefer explicit closing_image_query
    closing_query = _session_image_query(
        content if isinstance(content, dict) else {},
        project, session, 'closing',
    )
    has_img = _add_image_full_bleed(slide, sw, sh, closing_query,
                                    fallback_color=SLIDE_DEEP, variant=99)
    if has_img:
        _add_dark_overlay(slide, Inches(0), Inches(0), sw, sh,
                          alpha=82, top_hex='1E1B4B', bottom_hex='000000')
    else:
        _add_rect(slide, Inches(0), Inches(0), sw, sh,
                  fill=SLIDE_DEEP, shape=MSO_SHAPE.RECTANGLE)

    # Eyebrow
    _add_textbox(slide, Inches(1.5), Inches(2.4), Inches(10.3),
                 Inches(0.5), 'TAKE THIS WITH YOU', font_size=12, bold=True,
                 color=ACCENT_GOLD, align=PP_ALIGN.CENTER)
    # Decorative line
    _add_rect(slide, Inches(6.4), Inches(3.0), Inches(0.5), Inches(0.06),
              fill=ACCENT_GOLD, shape=MSO_SHAPE.RECTANGLE)
    # Quote
    box = slide.shapes.add_textbox(Inches(1), Inches(3.3),
                                   Inches(11.3), Inches(2.8))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _render_inline_pptx(p, f'"{closing}"', font_size=36,
                        color=WHITE, base_bold=True, italic=True)
    # Bottom signature
    _add_textbox(slide, Inches(1), Inches(6.5), Inches(11.3), Inches(0.4),
                 f'— Week {session.week.number}  ·  Block Period {session.number}',
                 font_size=12, italic=True,
                 color=PRGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.CENTER)


def _prefetch_session_images(content: dict, project: Project, session: Session):
    """Walk through `content` and concurrently prefetch every image the
    session PPT will need. Mirrors the variant assignment in slide builders."""
    requests_list = []

    # Cover (variant=0)
    requests_list.append({
        'query': _session_image_query(content, project, session, 'cover'),
        'variant': 0,
    })

    # Per-activity images
    for i, act in enumerate(content.get('activity_slides', []) or []):
        q = _activity_image_query(act, fallback=project.topic)
        # Activity intro (variant 1+i*3)
        requests_list.append({'query': q, 'variant': 1 + i * 3})
        # Content slide image (variant 2+i*3)
        requests_list.append({'query': q, 'variant': 2 + i * 3})
        # Media slide (variant 3+i*3)
        media = act.get('media_placeholder')
        if media and str(media).strip():
            requests_list.append({
                'query': _clean_caption(str(media), max_chars=200) or str(media),
                'variant': 3 + i * 3,
            })
        # Video slide (variant 4+i*3) — only if video_placeholder present OR
        # media mentions video
        video_text = act.get('video_placeholder')
        if not video_text and media and re.search(
                r'\bvideo|clip|footage|reel|short\b', str(media), re.IGNORECASE):
            video_text = str(media)
        if video_text and str(video_text).strip():
            requests_list.append({
                'query': _clean_caption(str(video_text), max_chars=200) or str(video_text),
                'variant': 4 + i * 3,
            })

    # Closing (variant 99)
    if content.get('closing_thought'):
        requests_list.append({
            'query': _session_image_query(content, project, session, 'closing'),
            'variant': 99,
        })

    # 8 concurrent workers — balances speed vs APIYI rate limits.
    prefetch_images(requests_list, max_workers=8)


def build_session_pptx(project: Project, session: Session, content: dict) -> io.BytesIO:
    """Build a kid-friendly, multi-slide Session PPT (~17-21 slides).

    Per-activity slide pattern:
      • Intro (full-bleed colour block with big number)
      • Content (cards with markdown-parsed text)
      • Discussion prompts (numbered cards)         — only if prompts present
      • Media placeholder (big dashed box)          — only if media specified
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    sw, sh = prs.slide_width, prs.slide_height
    blank_layout = prs.slide_layouts[6]

    # ── Prefetch all images concurrently — turns ~22 sequential DALL-E
    # ── calls (~11 min) into a single ~30-60 sec parallel batch.
    _prefetch_session_images(content, project, session)

    # Frame slides
    _slide_cover(prs, blank_layout, sw, sh, project, session, content)
    _slide_goals(prs, blank_layout, sw, sh, content)
    _slide_timeline(prs, blank_layout, sw, sh, content.get('timeline', []))

    # Per-activity slides — split across multiple slides each
    activity_slides = content.get('activity_slides', [])
    for i, act in enumerate(activity_slides):
        theme = ACTIVITY_THEMES[i % len(ACTIVITY_THEMES)]
        blocks  = [b for b in (act.get('content_blocks') or []) if str(b).strip()]
        prompts = [p for p in (act.get('prompts') or []) if str(p).strip()]
        media   = act.get('media_placeholder')

        # 1) Intro (full-bleed image hero)
        _slide_activity_intro(prs, blank_layout, sw, sh, i, act, theme, project)

        # 2) Content slide(s) — first one carries an image; split chunks of 3
        if blocks:
            CHUNK = 3
            for k in range(0, len(blocks), CHUNK):
                _slide_activity_content(prs, blank_layout, sw, sh, i, act, theme,
                                        blocks[k:k + CHUNK], project)

        # 3) Prompts
        if prompts:
            _slide_activity_prompts(prs, blank_layout, sw, sh, i, act, theme, prompts)

        # 4) Media image slide
        if media and str(media).strip():
            _slide_media_image(prs, blank_layout, sw, sh, i, act,
                               theme, str(media))

        # 5) Video placeholder slide (if AI provided one OR media mentions video)
        video_text = act.get('video_placeholder')
        if not video_text and media and re.search(
                r'\bvideo|clip|footage|reel|short\b', str(media), re.IGNORECASE):
            video_text = str(media)
        if video_text and str(video_text).strip():
            _slide_video_placeholder(prs, blank_layout, sw, sh, i, act,
                                     theme, str(video_text), project)

    # Closing frame slides
    _slide_reflection(prs, blank_layout, sw, sh, content.get('reflection', []))
    _slide_closing(prs, blank_layout, sw, sh, content, project, session)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
